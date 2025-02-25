from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import base64
import os
import json

class EtractDataPowerpoint:
    def __init__(self):
        self.json_pptx = None
    
    def extract_table_data(self, table):
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)
        return table_data
    
    def extract_image(self, shape, image_counter):
        image_data = {}
        try:
            image = shape.image
            image_bytes = image.blob
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            image_ext = image.ext
            image_data = {
                'image_id': f'image_{image_counter}',
                'content_type': f'image/{image_ext}',
                'data': image_base64
            }
        except:
            pass
        return image_data
    
    def extract_pptx_content(self, pptx_path):
        presentation = Presentation(pptx_path)
        presentation_data = {
            'slides': []
        }
        image_counter = 1
        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_data = {
                'slide_number': slide_number,
                'shapes': {
                    'text': [],
                    'tables': [],
                    'images': []
                }
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_data = []
                    for paragraph in shape.text_frame.paragraphs:
                        text_data.append({
                            'text': paragraph.text,
                            'level': paragraph.level
                        })
                    if text_data:
                        text_shape_data = {
                            'content': text_data,
                            'position': {
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height
                            }
                        }
                        slide_data['shapes']['text'].append(text_shape_data)
                elif shape.has_table:
                    table_data = {
                        'content': self.extract_table_data(shape.table),
                        'position': {
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        }
                    }
                    slide_data['shapes']['tables'].append(table_data)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = self.extract_image(shape, image_counter)
                    if image_data:
                        image_data['position'] = {
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        }
                        slide_data['shapes']['images'].append(image_data)
                        image_counter += 1
            presentation_data['slides'].append(slide_data)
        return presentation_data
    
    def save_to_json(self, data, output_path):
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
            
    def extract_data_from_pdf(self, pptx_path):
        self.json_pptx = json(self.extract_pptx_content(pptx_path))
        print(self.json_pptx)

# Sử dụng function
def main():
    pptx_path = '/home/nhatthuong/Documents/Thesis/Acne-detection-and-treatment-recommendations/backend/fastapi_all/ai/rag/convert_document/pptx_file/powerpoint_sample.pptx'  # Đường dẫn đến file PPTX
    output_path = 'output.json'  # Đường dẫn để lưu file JSON
    extract_pdf = EtractDataPowerpoint()
    try:
        # Trích xuất nội dung
        presentation_data = extract_pdf.extract_pptx_content(pptx_path)
        
        # Lưu vào file JSON
        extract_pdf.save_to_json(presentation_data, output_path)
        
        print(f"Đã xuất thành công dữ liệu ra file {output_path}")
    
    except Exception as e:
        print(f"Có lỗi xảy ra: {str(e)}")

if __name__ == "__main__":
    main()
